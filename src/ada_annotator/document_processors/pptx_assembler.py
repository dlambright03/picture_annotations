"""
PPTX document assembler.

Applies AI-generated alt-text to images in PowerPoint presentations
while preserving exact positions, sizes, rotations, and all visual
effects.
"""

from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from ada_annotator.document_processors.base_assembler import (
    DocumentAssembler
)
from ada_annotator.exceptions import ProcessingError
from ada_annotator.models import AltTextResult


class PPTXAssembler(DocumentAssembler):
    """
    Apply alt-text to PPTX presentations.

    Supports:
    - Picture shapes on slides
    - Position preservation (EMU precision)
    - Size preservation (width, height)
    - Rotation preservation
    - Visual effects preservation (shadow, glow, etc.)
    - Z-order preservation
    - Grouping preservation
    """

    def __init__(self, input_path: Path, output_path: Path):
        """
        Initialize PPTX assembler.

        Args:
            input_path: Path to input PPTX file.
            output_path: Path where output PPTX will be saved.

        Raises:
            FileNotFoundError: If input file does not exist.
            ValueError: If file is not PPTX format.
        """
        super().__init__(input_path, output_path)

        # Validate PPTX extension
        if input_path.suffix.lower() != ".pptx":
            raise ValueError(
                f"Not a PPTX file: {input_path}"
            )

        try:
            self.presentation = Presentation(str(input_path))
            self.logger.info(
                "pptx_loaded",
                slide_count=len(self.presentation.slides),
            )
        except Exception as e:
            raise ProcessingError(
                f"Failed to load PPTX presentation: {e}"
            ) from e

    def apply_alt_text(
        self,
        alt_text_results: List[AltTextResult]
    ) -> Dict[str, str]:
        """
        Apply alt-text to images in PPTX presentation.

        Finds images by slide and shape index, applies alt-text via
        XML manipulation. Preserves all shape properties including
        position, size, rotation, and visual effects.

        Args:
            alt_text_results: List of alt-text generation results.

        Returns:
            Dict[str, str]: Map of image_id to status message.

        Raises:
            ProcessingError: If critical errors occur.
        """
        status_map = {}

        self.logger.info(
            "applying_alt_text",
            total_results=len(alt_text_results),
        )

        for result in alt_text_results:
            try:
                status = self._apply_alt_text_to_image(result)
                status_map[result.image_id] = status

                self.logger.debug(
                    "alt_text_applied",
                    image_id=result.image_id,
                    status=status,
                    alt_text_length=len(result.alt_text),
                )
            except Exception as e:
                status_map[result.image_id] = f"failed: {str(e)}"
                self.logger.warning(
                    "alt_text_application_failed",
                    image_id=result.image_id,
                    error=str(e),
                )

        success_count = sum(
            1 for s in status_map.values() if s == "success"
        )
        self.logger.info(
            "alt_text_application_complete",
            success_count=success_count,
            total_count=len(alt_text_results),
        )

        return status_map

    def _apply_alt_text_to_image(
        self,
        result: AltTextResult
    ) -> str:
        """
        Apply alt-text to a single image.

        Args:
            result: Alt-text generation result.

        Returns:
            str: Status message ('success', 'skipped', or error).
        """
        # Extract slide and shape indices from image_id
        # Format: "slide{slide_idx}_shape{shape_idx}"
        try:
            parts = result.image_id.split("_")
            slide_idx = int(parts[0].replace("slide", ""))
            shape_idx = int(parts[1].replace("shape", ""))
        except (IndexError, ValueError):
            return "failed: invalid image_id format"

        # Validate slide index
        if slide_idx >= len(self.presentation.slides):
            return "failed: slide index out of range"

        slide = self.presentation.slides[slide_idx]

        # Find the picture shape
        shape = self._find_picture_shape(
            slide, slide_idx, shape_idx
        )

        if shape is None:
            return "failed: picture shape not found"

        # Apply alt-text to shape
        success = self._set_alt_text_on_shape(
            shape, result.alt_text
        )

        if success:
            return "success"
        else:
            return "failed: could not set alt-text"

    def _find_picture_shape(
        self,
        slide,
        slide_idx: int,
        shape_idx: int
    ) -> Optional:
        """
        Find picture shape on slide by index.

        Args:
            slide: The slide object.
            slide_idx: Zero-based slide index.
            shape_idx: Zero-based shape index on slide.

        Returns:
            Optional: Picture shape if found, None otherwise.
        """
        # Count picture shapes to find the correct one
        picture_count = 0

        for shape in slide.shapes:
            # Check if it's a picture shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if picture_count == shape_idx:
                    return shape
                picture_count += 1

        self.logger.debug(
            "picture_shape_not_found",
            slide_index=slide_idx,
            target_shape_index=shape_idx,
            picture_shapes_found=picture_count,
        )

        return None

    def _set_alt_text_on_shape(
        self,
        shape,
        alt_text: str
    ) -> bool:
        """
        Set alt-text on picture shape.

        Sets alt-text in both the shape name property and the
        XML cNvPr element for maximum compatibility.

        Args:
            shape: The picture shape object.
            alt_text: Alt-text to apply.

        Returns:
            bool: True if successful, False otherwise.
        """
        try:
            # Method 1: Set shape name (visible in PowerPoint UI)
            if hasattr(shape, 'name'):
                shape.name = alt_text

            # Method 2: Set XML attributes in cNvPr element
            # This is the standard location for alt-text
            if hasattr(shape, '_element'):
                # Find cNvPr element (non-visual picture properties)
                try:
                    # Path: p:sp/p:nvSpPr/p:cNvPr or
                    # p:pic/p:nvPicPr/p:cNvPr
                    nv_pr = shape._element.find(
                        './/' + qn('p:cNvPr')
                    )

                    if nv_pr is not None:
                        # Set both title and descr for compatibility
                        nv_pr.set('title', alt_text)
                        nv_pr.set('descr', alt_text)

                        self.logger.debug(
                            "xml_alt_text_set",
                            shape_name=shape.name if hasattr(
                                shape, 'name'
                            ) else "unknown",
                            alt_text_length=len(alt_text),
                        )

                except Exception as xml_error:
                    self.logger.warning(
                        "xml_alt_text_set_failed",
                        error=str(xml_error),
                    )

            return True

        except Exception as e:
            self.logger.error(
                "set_alt_text_failed",
                error=str(e),
            )
            return False

    def save_document(self) -> None:
        """
        Save the modified presentation to output path.

        Raises:
            ProcessingError: If save operation fails.
        """
        try:
            self.presentation.save(str(self.output_path))
            self.logger.info(
                "pptx_saved",
                output_path=str(self.output_path),
                file_size_bytes=self.output_path.stat().st_size,
            )
        except Exception as e:
            raise ProcessingError(
                f"Failed to save PPTX presentation: {e}"
            ) from e

    def get_document_format(self) -> str:
        """
        Get the document format identifier.

        Returns:
            str: 'PPTX'
        """
        return "PPTX"

    def validate_document(self) -> bool:
        """
        Validate input document can be processed.

        Returns:
            bool: True if document is valid, False otherwise.
        """
        # Call base validation
        if not super().validate_document():
            return False

        # Additional PPTX-specific validation
        try:
            return len(self.presentation.slides) > 0
        except Exception:
            return False

