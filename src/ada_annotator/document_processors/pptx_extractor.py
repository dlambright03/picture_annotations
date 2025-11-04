"""
PPTX document image extractor.

Extracts images from PowerPoint presentations with position metadata,
slide context, and existing alt-text. Preserves precise positioning
information for output generation.
"""

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ada_annotator.document_processors.base_extractor import DocumentExtractor
from ada_annotator.exceptions import ProcessingError
from ada_annotator.models import ImageMetadata


class PPTXExtractor(DocumentExtractor):
    """
    Extract images from PPTX presentations.

    Supports:
    - Picture shapes on slides
    - Precise position metadata (x, y, width, height in EMUs)
    - Slide-level context (titles)
    - Existing alt-text extraction
    """

    def __init__(self, document_path: Path):
        """
        Initialize PPTX extractor.

        Args:
            document_path: Path to PPTX file.

        Raises:
            FileNotFoundError: If file does not exist.
            ValueError: If file is not PPTX format.
        """
        super().__init__(document_path)

        # Validate PPTX extension
        if document_path.suffix.lower() != ".pptx":
            raise ValueError(f"Not a PPTX file: {document_path}")

        # Load presentation
        try:
            self.presentation = Presentation(str(document_path))
            self.logger.info(
                "pptx_loaded",
                document_path=str(document_path),
                slide_count=len(self.presentation.slides),
            )
        except Exception as e:
            raise ProcessingError(f"Failed to load PPTX: {e}") from e

    def get_document_format(self) -> str:
        """
        Get document format identifier.

        Returns:
            str: 'PPTX'
        """
        return "PPTX"

    def extract_images(self) -> list[ImageMetadata]:
        """
        Extract all images from PPTX presentation.

        Iterates through all slides and extracts picture shapes with
        precise position metadata and slide context.

        Returns:
            List[ImageMetadata]: List of extracted images.

        Raises:
            ProcessingError: If extraction fails.
        """
        self.logger.info(
            "extraction_started",
            document_format="PPTX",
            total_slides=len(self.presentation.slides),
        )

        images = []

        try:
            # Iterate through all slides
            for slide_idx, slide in enumerate(self.presentation.slides):
                # Extract slide title for context
                slide_title = self._extract_slide_title(slide)

                # Extract images from this slide
                slide_images = self._extract_images_from_slide(
                    slide, slide_idx, slide_title
                )
                images.extend(slide_images)

            self.logger.info(
                "extraction_completed",
                total_images=len(images),
                total_slides=len(self.presentation.slides),
            )

            return images

        except Exception as e:
            self.logger.error(
                "extraction_failed",
                error=str(e),
            )
            raise ProcessingError(f"Failed to extract images: {e}") from e

    def _extract_slide_title(self, slide) -> str | None:
        """
        Extract title from slide for context.

        Args:
            slide: The slide object.

        Returns:
            Optional[str]: Slide title if present, None otherwise.
        """
        try:
            # Check if slide has title placeholder
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    return title_text
            return None
        except Exception:
            return None

    def _extract_images_from_slide(
        self, slide, slide_idx: int, slide_title: str | None
    ) -> list[ImageMetadata]:
        """
        Extract all images from a single slide.

        Args:
            slide: The slide object.
            slide_idx: Zero-based slide index.
            slide_title: Slide title for context.

        Returns:
            List[ImageMetadata]: List of images from this slide.
        """
        images = []

        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Check if shape is a picture
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                # Extract image data
                image_metadata = self._extract_image_from_shape(
                    shape, slide_idx, shape_idx, slide_title
                )

                if image_metadata:
                    images.append(image_metadata)

            except Exception as e:
                self.logger.warning(
                    "shape_extraction_failed",
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    error=str(e),
                )
                continue

        self.logger.debug(
            "slide_extraction_completed",
            slide_index=slide_idx,
            images_found=len(images),
            slide_title=slide_title,
        )

        return images

    def _extract_image_from_shape(
        self, shape, slide_idx: int, shape_idx: int, slide_title: str | None
    ) -> ImageMetadata | None:
        """
        Extract image metadata from a picture shape.

        Args:
            shape: The picture shape object.
            slide_idx: Zero-based slide index.
            shape_idx: Zero-based shape index on slide.
            slide_title: Slide title for context.

        Returns:
            Optional[ImageMetadata]: Image metadata if successful.
        """
        try:
            # Get image binary data
            image_bytes = shape.image.blob

            # Get image format from content type
            content_type = shape.image.content_type
            format_str = content_type.split("/")[-1].upper()

            # Normalize format names
            if format_str in ["JPG", "JPEG"]:
                format_str = "JPEG"

            # Load with PIL to get dimensions
            img = Image.open(BytesIO(image_bytes))

            # Extract existing alt-text
            alt_text = self._extract_alt_text_from_shape(shape)

            # Create position metadata (EMUs - English Metric Units)
            # 1 inch = 914400 EMUs
            position = {
                "slide_index": slide_idx,
                "shape_index": shape_idx,
                "left_emu": shape.left,
                "top_emu": shape.top,
                "width_emu": shape.width,
                "height_emu": shape.height,
                "slide_title": slide_title,
            }

            # Create unique image ID
            image_id = f"slide{slide_idx}_shape{shape_idx}"

            # Create ImageMetadata with binary data
            metadata = ImageMetadata(
                image_id=image_id,
                filename=f"{image_id}.{format_str.lower()}",
                format=format_str,
                size_bytes=len(image_bytes),
                width_pixels=img.width,
                height_pixels=img.height,
                page_number=slide_idx + 1,  # 1-based slide number
                position=position,
                existing_alt_text=alt_text,
                image_data=image_bytes,  # Store binary data for processing
            )

            self.logger.debug(
                "image_extracted",
                image_id=image_id,
                slide_index=slide_idx,
                shape_index=shape_idx,
                format=format_str,
                has_alt_text=bool(alt_text),
            )

            return metadata

        except Exception as e:
            self.logger.warning(
                "image_extraction_failed",
                slide_index=slide_idx,
                shape_index=shape_idx,
                error=str(e),
            )
            return None

    def _extract_alt_text_from_shape(self, shape) -> str | None:
        """
        Extract existing alt-text from picture shape.

        Args:
            shape: The picture shape object.

        Returns:
            Optional[str]: Alt-text if present, None otherwise.
        """
        try:
            # Check shape name (sometimes contains alt-text)
            if hasattr(shape, "name") and shape.name:
                name = shape.name.strip()
                # Ignore default names like "Picture 1", "Image 2"
                if not (
                    name.startswith("Picture") or name.startswith("Image")
                ):
                    return name

            # Check for title/description in shape element
            if hasattr(shape, "_element"):
                # Try to find cNvPr (non-visual properties)
                nvPr = shape._element.xpath(
                    ".//p:cNvPr",
                    namespaces={
                        "p": (
                            "http://schemas.openxmlformats.org/"
                            "presentationml/2006/main"
                        )
                    },
                )

                for prop in nvPr:
                    # Check for title attribute
                    title = prop.get("title")
                    if title and title.strip():
                        return title.strip()

                    # Check for descr attribute
                    descr = prop.get("descr")
                    if descr and descr.strip():
                        return descr.strip()

            return None

        except Exception:
            return None
