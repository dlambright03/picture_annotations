"""
Context extraction for images from DOCX and PPTX documents.

Implements a hierarchical context extraction system with 5 levels:
1. External context (from user-provided file)
2. Document context (title, subject, author metadata)
3. Section context (nearest heading before image)
4. Page context (slide title for PPTX)
5. Local context (surrounding paragraphs)
"""

from pathlib import Path
from typing import List, Optional, Union

from docx import Document as DocxDocument
from pptx import Presentation

from ada_annotator.models import ContextData, ImageMetadata
from ada_annotator.utils.logging import get_logger


class ContextExtractor:
    """
    Extract hierarchical context for images in documents.

    Supports both DOCX and PPTX formats with document-specific
    context extraction strategies.
    """

    def __init__(
        self,
        document_path: Path,
        external_context_path: Optional[Path] = None
    ):
        """
        Initialize context extractor.

        Args:
            document_path: Path to DOCX or PPTX file.
            external_context_path: Optional path to external context
                                   file (.txt or .md).

        Raises:
            ValueError: If file format is unsupported.
        """
        self.document_path = document_path
        self.external_context_path = external_context_path
        self.logger = get_logger(__name__)

        # Load external context if provided
        self.external_context = None
        if external_context_path:
            self.external_context = self._load_external_context(
                external_context_path
            )

        # Determine document type and load
        suffix = document_path.suffix.lower()
        if suffix == ".docx":
            self.docx_document = DocxDocument(str(document_path))
            self.pptx_document = None
            self.document_type = "DOCX"
        elif suffix == ".pptx":
            self.docx_document = None
            self.pptx_document = Presentation(str(document_path))
            self.document_type = "PPTX"
        else:
            raise ValueError(
                f"Unsupported document format: {suffix}"
            )

        self.logger.info(
            "context_extractor_initialized",
            document_type=self.document_type,
            has_external_context=bool(self.external_context),
        )

    def extract_context_for_image(
        self,
        image_metadata: ImageMetadata
    ) -> ContextData:
        """
        Extract complete hierarchical context for an image.

        Args:
            image_metadata: Metadata for the image.

        Returns:
            ContextData: Complete context hierarchy.
        """
        # Extract document-level context
        document_context = self._extract_document_context()

        # Extract section context
        section_context = self._extract_section_context(
            image_metadata
        )

        # Extract page context (PPTX only)
        page_context = self._extract_page_context(image_metadata)

        # Extract local context
        local_context = self._extract_local_context(
            image_metadata
        )

        context_data = ContextData(
            external_context=self.external_context,
            document_context=document_context,
            section_context=section_context,
            page_context=page_context,
            local_context=local_context,
        )

        self.logger.debug(
            "context_extracted",
            image_id=image_metadata.image_id,
            has_external=bool(self.external_context),
            has_section=bool(section_context),
            has_page=bool(page_context),
            local_length=len(local_context),
        )

        return context_data

    def _load_external_context(
        self, context_path: Path
    ) -> Optional[str]:
        """
        Load external context from text or markdown file.

        Args:
            context_path: Path to context file (.txt or .md).

        Returns:
            Optional[str]: Context text, or None if file not found.
        """
        try:
            # Validate file extension
            if context_path.suffix.lower() not in [".txt", ".md"]:
                self.logger.warning(
                    "unsupported_context_file_format",
                    file_path=str(context_path),
                    format=context_path.suffix,
                )
                return None

            # Read file
            with open(context_path, "r", encoding="utf-8") as f:
                content = f.read().strip()

            # Validate content length (max 10000 chars)
            if len(content) > 10000:
                self.logger.warning(
                    "external_context_too_long",
                    length=len(content),
                    max_length=10000,
                )
                content = content[:10000] + "..."

            self.logger.info(
                "external_context_loaded",
                file_path=str(context_path),
                length=len(content),
            )

            return content

        except FileNotFoundError:
            self.logger.warning(
                "external_context_file_not_found",
                file_path=str(context_path),
            )
            return None
        except Exception as e:
            self.logger.warning(
                "external_context_load_failed",
                file_path=str(context_path),
                error=str(e),
            )
            return None

    def _extract_document_context(self) -> str:
        """
        Extract document-level metadata context.

        For DOCX: core properties (title, subject, author)
        For PPTX: core properties (title, subject, author)

        Returns:
            str: Formatted document context string.
        """
        try:
            parts = []

            if self.document_type == "DOCX":
                props = self.docx_document.core_properties  # type: ignore
            else:  # PPTX
                props = self.pptx_document.core_properties  # type: ignore

            # Extract title
            if props.title:
                parts.append(f"Title: {props.title}")

            # Extract subject
            if props.subject:
                parts.append(f"Subject: {props.subject}")

            # Extract author
            if props.author:
                parts.append(f"Author: {props.author}")

            # Default if no metadata
            if not parts:
                parts.append(
                    f"{self.document_type} document "
                    f"({self.document_path.name})"
                )

            context = ", ".join(parts)

            self.logger.debug(
                "document_context_extracted",
                context=context,
            )

            return context

        except Exception as e:
            self.logger.warning(
                "document_context_extraction_failed",
                error=str(e),
            )
            return f"{self.document_type} document"

    def _extract_section_context(
        self,
        image_metadata: ImageMetadata
    ) -> Optional[str]:
        """
        Extract section context (nearest heading before image).

        For DOCX: Searches backwards from paragraph index
        For PPTX: Uses slide title (if available)

        Args:
            image_metadata: Metadata for the image.

        Returns:
            Optional[str]: Section heading text, or None.
        """
        if self.document_type == "DOCX":
            return self._extract_docx_section_context(
                image_metadata
            )
        else:  # PPTX
            return self._extract_pptx_section_context(
                image_metadata
            )

    def _extract_docx_section_context(
        self,
        image_metadata: ImageMetadata
    ) -> Optional[str]:
        """
        Find nearest heading before image in DOCX document.

        Args:
            image_metadata: Metadata for the image.

        Returns:
            Optional[str]: Heading text, or None if no heading found.
        """
        try:
            # Get paragraph index from position metadata
            para_index = image_metadata.position.get(
                "paragraph_index", 0
            )

            # Search backwards from image paragraph
            for i in range(para_index - 1, -1, -1):
                para = self.docx_document.paragraphs[i]  # type: ignore

                # Check if paragraph is a heading
                if para.style and para.style.name and \
                   para.style.name.startswith("Heading"):
                    heading_text = para.text.strip()
                    if heading_text:
                        self.logger.debug(
                            "section_context_found",
                            image_id=image_metadata.image_id,
                            heading=heading_text,
                            heading_level=para.style.name,
                        )
                        return heading_text

            # No heading found
            return None

        except Exception as e:
            self.logger.warning(
                "section_context_extraction_failed",
                image_id=image_metadata.image_id,
                error=str(e),
            )
            return None

    def _extract_pptx_section_context(
        self,
        image_metadata: ImageMetadata
    ) -> Optional[str]:
        """
        Get slide title as section context for PPTX.

        Args:
            image_metadata: Metadata for the image.

        Returns:
            Optional[str]: Slide title, or None.
        """
        try:
            # Check if slide_title is in position metadata
            slide_title = image_metadata.position.get("slide_title")
            return slide_title

        except Exception:
            return None

    def _extract_page_context(
        self,
        image_metadata: ImageMetadata
    ) -> Optional[str]:
        """
        Extract page/slide context.

        For DOCX: None (no page concept)
        For PPTX: Slide title

        Args:
            image_metadata: Metadata for the image.

        Returns:
            Optional[str]: Page context, or None.
        """
        if self.document_type == "PPTX":
            # Get slide title from position metadata
            slide_title = image_metadata.position.get("slide_title")
            if slide_title:
                return f"Slide: {slide_title}"

        return None

    def _extract_local_context(
        self,
        image_metadata: ImageMetadata,
        paragraphs_before: int = 2,
        paragraphs_after: int = 2
    ) -> str:
        """
        Extract local context (surrounding paragraphs).

        Args:
            image_metadata: Metadata for the image.
            paragraphs_before: Number of paragraphs before image.
            paragraphs_after: Number of paragraphs after image.

        Returns:
            str: Concatenated local context.
        """
        if self.document_type == "DOCX":
            return self._extract_docx_local_context(
                image_metadata,
                paragraphs_before,
                paragraphs_after
            )
        else:  # PPTX
            return self._extract_pptx_local_context(
                image_metadata
            )

    def _extract_docx_local_context(
        self,
        image_metadata: ImageMetadata,
        paragraphs_before: int,
        paragraphs_after: int
    ) -> str:
        """
        Extract surrounding paragraphs in DOCX.

        Args:
            image_metadata: Metadata for the image.
            paragraphs_before: Number of paragraphs before image.
            paragraphs_after: Number of paragraphs after image.

        Returns:
            str: Concatenated paragraph text.
        """
        try:
            # Get paragraph index from position metadata
            para_index = image_metadata.position.get(
                "paragraph_index", 0
            )

            # Calculate range
            start_idx = max(0, para_index - paragraphs_before)
            end_idx = min(
                len(self.docx_document.paragraphs),  # type: ignore
                para_index + paragraphs_after + 1
            )

            # Collect paragraph text
            paragraphs = []
            for i in range(start_idx, end_idx):
                # Skip the image paragraph itself
                if i == para_index:
                    continue

                para = self.docx_document.paragraphs[i]  # type: ignore
                text = para.text.strip()

                # Skip empty paragraphs
                if text:
                    paragraphs.append(text)

            # Join with separator
            context = " ".join(paragraphs)

            # If no context found, provide default
            if not context:
                context = "No surrounding text available."

            self.logger.debug(
                "local_context_extracted",
                image_id=image_metadata.image_id,
                paragraphs_collected=len(paragraphs),
                context_length=len(context),
            )

            return context

        except Exception as e:
            self.logger.warning(
                "local_context_extraction_failed",
                image_id=image_metadata.image_id,
                error=str(e),
            )
            return "No surrounding text available."

    def _extract_pptx_local_context(
        self,
        image_metadata: ImageMetadata
    ) -> str:
        """
        Extract text from slide for PPTX local context.

        Args:
            image_metadata: Metadata for the image.

        Returns:
            str: Slide text content.
        """
        try:
            # Get slide index from position metadata
            slide_idx = image_metadata.position.get("slide_index", 0)

            # Get the slide
            slide = self.pptx_document.slides[slide_idx]  # type: ignore

            # Collect all text from shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()  # type: ignore
                    if text:
                        text_parts.append(text)

            # Join with separator
            context = " ".join(text_parts)

            # If no context found, provide default
            if not context:
                context = "No text content on slide."

            self.logger.debug(
                "local_context_extracted",
                image_id=image_metadata.image_id,
                shapes_collected=len(text_parts),
                context_length=len(context),
            )

            return context

        except Exception as e:
            self.logger.warning(
                "local_context_extraction_failed",
                image_id=image_metadata.image_id,
                error=str(e),
            )
            return "No text content on slide."
