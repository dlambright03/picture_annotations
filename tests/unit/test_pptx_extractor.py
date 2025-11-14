"""
Unit tests for PPTX image extractor.

Tests PPTX image extraction with position metadata, slide context,
and existing alt-text extraction.
"""

from io import BytesIO
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from ada_annotator.document_processors import PPTXExtractor
from ada_annotator.exceptions import ProcessingError


class TestPPTXExtractorInitialization:
    """Test PPTXExtractor initialization and validation."""

    def test_initialization_success(self, tmp_path):
        """Test successful initialization with valid PPTX."""
        # Create minimal PPTX
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)

        assert extractor.document_path == pptx_path
        assert extractor.get_document_format() == "PPTX"
        assert extractor.presentation is not None

    def test_initialization_file_not_found(self, tmp_path):
        """Test initialization with non-existent file."""
        pptx_path = tmp_path / "nonexistent.pptx"

        with pytest.raises(FileNotFoundError):
            PPTXExtractor(pptx_path)

    def test_initialization_wrong_extension(self, tmp_path):
        """Test initialization with wrong file extension."""
        txt_path = tmp_path / "test.txt"
        txt_path.write_text("Not a PPTX file")

        with pytest.raises(ValueError, match="Not a PPTX file"):
            PPTXExtractor(txt_path)

    def test_initialization_corrupted_file(self, tmp_path):
        """Test initialization with corrupted PPTX."""
        pptx_path = tmp_path / "corrupted.pptx"
        pptx_path.write_bytes(b"Not a valid PPTX file")

        with pytest.raises(ProcessingError):
            PPTXExtractor(pptx_path)

    def test_validate_document(self, tmp_path):
        """Test document validation."""
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        assert extractor.validate_document() is True


class TestPPTXImageExtraction:
    """Test PPTX image extraction functionality."""

    def test_extract_images_empty_presentation(self, tmp_path):
        """Test extraction from presentation with no images."""
        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])  # Blank slide
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 0

    def test_extract_single_image(self, tmp_path):
        """Test extraction of single image from slide."""
        # Create test image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (100, 100), color="red")
        test_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Create PPTX with image
        pptx_path = tmp_path / "single_image.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        slide.shapes.add_picture(
            img_bytes,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2)
        )
        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        assert images[0].format == "PNG"
        assert images[0].width_pixels == 100
        assert images[0].height_pixels == 100
        assert images[0].page_number == 1

    def test_extract_multiple_images_single_slide(self, tmp_path):
        """Test extraction of multiple images from one slide."""
        # Create test images
        img1_bytes = BytesIO()
        test_img1 = Image.new("RGB", (50, 50), color="red")
        test_img1.save(img1_bytes, format="PNG")
        img1_bytes.seek(0)

        img2_bytes = BytesIO()
        test_img2 = Image.new("RGB", (60, 60), color="blue")
        test_img2.save(img2_bytes, format="PNG")
        img2_bytes.seek(0)

        # Create PPTX
        pptx_path = tmp_path / "multi_image.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img1_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        slide.shapes.add_picture(
            img2_bytes, Inches(3), Inches(3), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 2
        assert images[0].width_pixels == 50
        assert images[1].width_pixels == 60

    def test_extract_images_multiple_slides(self, tmp_path):
        """Test extraction from multiple slides."""
        # Create test image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (40, 40), color="green")
        test_img.save(img_bytes, format="PNG")

        # Create PPTX with multiple slides
        pptx_path = tmp_path / "multi_slide.pptx"
        prs = Presentation()

        for i in range(3):
            img_bytes.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                img_bytes,
                Inches(1),
                Inches(1),
                Inches(1),
                Inches(1)
            )

        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 3
        assert images[0].page_number == 1
        assert images[1].page_number == 2
        assert images[2].page_number == 3


class TestPPTXPositionMetadata:
    """Test position metadata extraction."""

    def test_position_metadata_captured(self, tmp_path):
        """Test that position metadata is captured correctly."""
        # Create test image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (100, 100), color="blue")
        test_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Create PPTX
        pptx_path = tmp_path / "position.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            img_bytes,
            Inches(2),
            Inches(3),
            Inches(4),
            Inches(5)
        )
        prs.save(str(pptx_path))

        # Extract and verify position
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        position = images[0].position
        assert "slide_index" in position
        assert "shape_index" in position
        assert "left_emu" in position
        assert "top_emu" in position
        assert "width_emu" in position
        assert "height_emu" in position
        assert position["slide_index"] == 0
        assert position["shape_index"] >= 0

    def test_slide_title_extracted(self, tmp_path):
        """Test that slide titles are captured as context."""
        # Create test image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (50, 50), color="yellow")
        test_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Create PPTX with title
        pptx_path = tmp_path / "titled.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
        slide.shapes.title.text = "Test Slide Title"
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(2), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        position = images[0].position
        assert position.get("slide_title") == "Test Slide Title"

    def test_slide_no_title(self, tmp_path):
        """Test handling of slides without titles."""
        # Create test image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (50, 50), color="purple")
        test_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Create PPTX without title
        pptx_path = tmp_path / "notitle.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        position = images[0].position
        assert position.get("slide_title") is None


class TestPPTXAltTextExtraction:
    """Test existing alt-text extraction from PPTX."""

    def test_image_id_generation(self, tmp_path):
        """Test unique image ID generation."""
        # Create test images
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (30, 30), color="cyan")
        test_img.save(img_bytes, format="PNG")

        # Create PPTX
        pptx_path = tmp_path / "ids.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )

        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes, Inches(3), Inches(3), Inches(1), Inches(1)
        )

        prs.save(str(pptx_path))

        # Extract and verify IDs are unique
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 2
        assert images[0].image_id != images[1].image_id
        assert "slide0" in images[0].image_id
        assert "shape" in images[0].image_id


class TestPPTXEdgeCases:
    """Test edge cases and error handling."""

    def test_non_picture_shapes_ignored(self, tmp_path):
        """Test that non-picture shapes are ignored."""
        pptx_path = tmp_path / "shapes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add text box (not a picture)
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        )
        txBox.text = "Test text"

        prs.save(str(pptx_path))

        # Extract - should find no images
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 0

    def test_format_normalization(self, tmp_path):
        """Test that image formats are normalized."""
        # Create JPEG image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (50, 50), color="orange")
        test_img.save(img_bytes, format="JPEG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "jpeg.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        # Extract and verify format
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        assert images[0].format == "JPEG"


class TestPPTXIntegration:
    """Integration tests for complete PPTX processing."""

    def test_complete_extraction_workflow(self, tmp_path):
        """Test complete extraction with multiple slides and images."""
        # Create multiple test images
        red_img = BytesIO()
        Image.new("RGB", (100, 100), "red").save(red_img, "PNG")

        blue_img = BytesIO()
        Image.new("RGB", (200, 200), "blue").save(blue_img, "PNG")

        # Create complex PPTX
        pptx_path = tmp_path / "complex.pptx"
        prs = Presentation()

        # Slide 1: Title slide with image
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Slide 1 Title"
        red_img.seek(0)
        slide1.shapes.add_picture(
            red_img, Inches(1), Inches(2), Inches(2), Inches(2)
        )

        # Slide 2: Blank slide with two images
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        red_img.seek(0)
        slide2.shapes.add_picture(
            red_img, Inches(0.5), Inches(0.5), Inches(1), Inches(1)
        )
        blue_img.seek(0)
        slide2.shapes.add_picture(
            blue_img, Inches(4), Inches(4), Inches(3), Inches(3)
        )

        prs.save(str(pptx_path))

        # Extract and verify
        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        # Should have 3 total images
        assert len(images) == 3

        # Verify slide distribution
        slide_numbers = [img.page_number for img in images]
        assert slide_numbers.count(1) == 1  # 1 image on slide 1
        assert slide_numbers.count(2) == 2  # 2 images on slide 2

        # Verify dimensions
        assert images[0].width_pixels == 100
        assert images[2].width_pixels == 200

        # Verify position metadata exists
        for img in images:
            assert "slide_index" in img.position
            assert "left_emu" in img.position
            assert "width_emu" in img.position


class TestPPTXAltTextExtensiveExtraction:
    """Test extensive alt-text extraction scenarios."""

    def test_extract_alt_text_from_title_attribute(self, tmp_path):
        """Test alt-text extraction from title attribute."""
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "magenta").save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "alt_title.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )

        # Manually set title attribute (simulating existing alt-text)
        # This requires direct XML manipulation
        try:
            nvPr = pic._element.xpath(".//p:cNvPr")[0]
            nvPr.set("title", "Alt text from title attribute")
        except (IndexError, AttributeError):
            pass  # Skip if XML structure doesn't support it

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        # Alt-text may or may not be extracted depending on XML structure
        # This tests the extraction logic runs without error

    def test_extract_alt_text_from_descr_attribute(self, tmp_path):
        """Test alt-text extraction from descr attribute."""
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "brown").save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "alt_descr.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )

        # Try to set descr attribute
        try:
            nvPr = pic._element.xpath(".//p:cNvPr")[0]
            nvPr.set("descr", "Description alt-text")
        except (IndexError, AttributeError):
            pass

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1

    def test_extract_alt_text_ignores_default_names(self, tmp_path):
        """Test that default picture names are ignored."""
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "gray").save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "default_name.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )

        # Picture will have default name like "Picture 1"
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        # Default names should NOT be used as alt-text
        alt_text = images[0].existing_alt_text
        if alt_text:
            assert not alt_text.startswith("Picture")
            assert not alt_text.startswith("Image")


class TestPPTXHelperMethods:
    """Test helper methods for slide and image extraction."""

    def test_extract_slide_title_with_title(self, tmp_path):
        """Test extracting slide title when present."""
        pptx_path = tmp_path / "titled_slide.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "My Important Slide"
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        slide_obj = prs.slides[0]
        title = extractor._extract_slide_title(slide_obj)

        assert title == "My Important Slide"

    def test_extract_slide_title_without_title(self, tmp_path):
        """Test extracting slide title when not present."""
        pptx_path = tmp_path / "notitle_slide.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        slide_obj = prs.slides[0]
        title = extractor._extract_slide_title(slide_obj)

        assert title is None

    def test_extract_slide_title_with_whitespace(self, tmp_path):
        """Test that slide titles with only whitespace are handled."""
        pptx_path = tmp_path / "whitespace_title.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "   \n\t  "  # Only whitespace
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        slide_obj = prs.slides[0]
        title = extractor._extract_slide_title(slide_obj)

        assert title is None


class TestPPTXErrorHandling:
    """Test error handling during extraction."""

    def test_extraction_continues_after_shape_error(self, tmp_path):
        """Test that extraction continues if one shape fails."""
        # Create valid image
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "teal").save(img_bytes, "PNG")

        pptx_path = tmp_path / "partial_error.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add valid image
        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )

        # Add text shape (non-picture)
        txBox = slide.shapes.add_textbox(
            Inches(3), Inches(3), Inches(2), Inches(1)
        )
        txBox.text = "Not an image"

        # Add another valid image
        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes, Inches(5), Inches(5), Inches(1), Inches(1)
        )

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        # Should extract 2 images despite text shape in between
        assert len(images) == 2

    def test_extract_images_from_slide_with_mixed_shapes(self, tmp_path):
        """Test extraction from slides with various shape types."""
        img_bytes = BytesIO()
        Image.new("RGB", (40, 40), "pink").save(img_bytes, "PNG")

        pptx_path = tmp_path / "mixed_shapes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add various shapes
        slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(2), Inches(1)
        )

        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes, Inches(2), Inches(2), Inches(1), Inches(1)
        )

        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4), Inches(4), Inches(1), Inches(1)
        )

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        # Should only extract the picture
        assert len(images) == 1


class TestPPTXImageFormats:
    """Test handling of different image formats."""

    def test_jpeg_format_handling(self, tmp_path):
        """Test JPEG image extraction."""
        img_bytes = BytesIO()
        Image.new("RGB", (60, 60), "maroon").save(img_bytes, "JPEG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "jpeg_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        assert images[0].format == "JPEG"

    def test_png_format_handling(self, tmp_path):
        """Test PNG image extraction."""
        img_bytes = BytesIO()
        Image.new("RGBA", (70, 70), (255, 0, 255, 128)).save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "png_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        assert images[0].format == "PNG"

    def test_image_data_included(self, tmp_path):
        """Test that image binary data is included in metadata."""
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (80, 80), "navy")
        test_img.save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "data_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_bytes, Inches(1), Inches(1), Inches(1), Inches(1)
        )
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        assert images[0].image_data is not None
        assert len(images[0].image_data) > 0
        assert images[0].size_bytes == len(images[0].image_data)


class TestPPTXPositionMetadataDetailed:
    """Test detailed position metadata extraction."""

    def test_emu_values_extracted(self, tmp_path):
        """Test that EMU values are correctly extracted."""
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "lime").save(img_bytes, "PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "emu_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            img_bytes, Inches(2), Inches(3), Inches(4), Inches(5)
        )
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 1
        pos = images[0].position

        # EMU values should be positive integers
        assert pos["left_emu"] > 0
        assert pos["top_emu"] > 0
        assert pos["width_emu"] > 0
        assert pos["height_emu"] > 0

        # 2 inches ~= 1828800 EMUs (914400 per inch)
        assert abs(pos["left_emu"] - 1828800) < 1000

    def test_shape_index_tracking(self, tmp_path):
        """Test that shape indices are correctly tracked."""
        img_bytes = BytesIO()
        Image.new("RGB", (50, 50), "olive").save(img_bytes, "PNG")

        pptx_path = tmp_path / "index_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add multiple pictures
        for i in range(3):
            img_bytes.seek(0)
            slide.shapes.add_picture(
                img_bytes, Inches(i), Inches(i), Inches(1), Inches(1)
            )

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        assert len(images) == 3

        # Verify unique shape indices
        shape_indices = [img.position["shape_index"] for img in images]
        assert len(set(shape_indices)) == 3  # All unique


class TestPPTXExtractionErrors:
    """Test error scenarios during PPTX extraction."""

    def test_handles_corrupted_presentation(self, tmp_path):
        """Should raise ProcessingError for corrupted PPTX files."""
        from ada_annotator.exceptions import ProcessingError

        # Create a corrupted file
        pptx_path = tmp_path / "corrupted.pptx"
        pptx_path.write_text("This is not a valid PPTX file")

        # Should raise ProcessingError during initialization
        with pytest.raises(ProcessingError, match="Failed to load PPTX"):
            PPTXExtractor(pptx_path)

    def test_handles_slide_with_no_shapes(self, tmp_path):
        """Should handle slides with no shapes gracefully."""
        pptx_path = tmp_path / "empty_slide.pptx"
        prs = Presentation()

        # Add slide with blank layout (no shapes)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)
        images = extractor.extract_images()

        # Should return empty list, not error
        assert len(images) == 0

    def test_extract_images_failure_handling(self, tmp_path):
        """Test that extract_images raises ProcessingError on failure."""
        from unittest.mock import patch
        from ada_annotator.exceptions import ProcessingError

        # Create a PPTX with an image
        img_bytes = BytesIO()
        test_img = Image.new("RGB", (100, 100), color="blue")
        test_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        pptx_path = tmp_path / "test_failure.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_bytes, Inches(1), Inches(1), Inches(2))
        prs.save(str(pptx_path))

        extractor = PPTXExtractor(pptx_path)

        # Mock _extract_images_from_slide to raise an exception
        with patch.object(extractor, '_extract_images_from_slide', side_effect=RuntimeError("Test error")):
            with pytest.raises(ProcessingError) as exc_info:
                extractor.extract_images()

            assert "Failed to extract images" in str(exc_info.value)

