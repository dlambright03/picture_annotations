"""
Unit tests for context extraction module.

Tests all 5 levels of context hierarchy and merging logic.
"""

from pathlib import Path
from unittest.mock import Mock, patch

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from ada_annotator.models import ContextData, ImageMetadata
from ada_annotator.utils import ContextExtractor


class TestContextExtractorInit:
    """Test ContextExtractor initialization."""

    def test_init_docx_without_external_context(
        self, sample_docx_path
    ):
        """Test initialization with DOCX and no external context."""
        extractor = ContextExtractor(sample_docx_path)

        assert extractor.document_type == "DOCX"
        assert extractor.docx_document is not None
        assert extractor.pptx_document is None
        assert extractor.external_context is None

    def test_init_pptx_without_external_context(
        self, sample_pptx_path
    ):
        """Test initialization with PPTX and no external context."""
        extractor = ContextExtractor(sample_pptx_path)

        assert extractor.document_type == "PPTX"
        assert extractor.pptx_document is not None
        assert extractor.docx_document is None
        assert extractor.external_context is None

    def test_init_with_external_context(
        self, sample_docx_path, sample_context_path
    ):
        """Test initialization with external context file."""
        extractor = ContextExtractor(
            sample_docx_path,
            external_context_path=sample_context_path
        )

        assert extractor.external_context is not None
        assert len(extractor.external_context) > 0

    def test_init_pdf_format(self, tmp_path):
        """Test initialization with PDF format (now supported)."""
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_text("fake pdf")

        # PDF is now supported, should not raise
        extractor = ContextExtractor(pdf_path)
        assert extractor is not None
        assert extractor.document_type == "PDF"


class TestExternalContextLoading:
    """Test external context file loading."""

    def test_load_txt_file(self, sample_context_path, sample_docx_path):
        """Test loading .txt external context file."""
        extractor = ContextExtractor(
            sample_docx_path,
            external_context_path=sample_context_path
        )

        assert extractor.external_context is not None
        assert "sample external context" in \
            extractor.external_context

    def test_load_md_file(self, tmp_path):
        """Test loading .md external context file."""
        md_path = tmp_path / "context.md"
        md_path.write_text("# Markdown Context\n\nTest content")

        # Create a dummy docx for the extractor
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.save(str(docx_path))

        extractor = ContextExtractor(
            docx_path,
            external_context_path=md_path
        )

        assert extractor.external_context is not None
        assert "Markdown Context" in extractor.external_context

    def test_load_nonexistent_file(self, sample_docx_path):
        """Test loading non-existent external context file."""
        fake_path = Path("nonexistent.txt")

        extractor = ContextExtractor(
            sample_docx_path,
            external_context_path=fake_path
        )

        assert extractor.external_context is None

    def test_load_unsupported_format(
        self, sample_docx_path, tmp_path
    ):
        """Test loading unsupported external context format."""
        json_path = tmp_path / "context.json"
        json_path.write_text('{"key": "value"}')

        extractor = ContextExtractor(
            sample_docx_path,
            external_context_path=json_path
        )

        assert extractor.external_context is None

    def test_truncate_long_context(self, sample_docx_path, tmp_path):
        """Test truncation of very long external context."""
        long_path = tmp_path / "long_context.txt"
        long_content = "x" * 15000  # Exceeds 10000 char limit
        long_path.write_text(long_content)

        extractor = ContextExtractor(
            sample_docx_path,
            external_context_path=long_path
        )

        assert extractor.external_context is not None
        assert len(extractor.external_context) <= 10003  # 10000 + "..."


class TestDocumentContextExtraction:
    """Test document-level context extraction."""

    def test_docx_with_metadata(self, tmp_path):
        """Test DOCX with complete metadata."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.core_properties.title = "Test Document"
        doc.core_properties.subject = "Unit Testing"
        doc.core_properties.author = "Test Author"
        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)
        context = extractor._extract_document_context()

        assert "Title: Test Document" in context
        assert "Subject: Unit Testing" in context
        assert "Author: Test Author" in context

    def test_docx_without_metadata(self, tmp_path):
        """Test DOCX without metadata."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)
        context = extractor._extract_document_context()

        # python-docx adds default author, so check for filename
        assert "test.docx" in context or "Author:" in context

    def test_pptx_with_metadata(self, tmp_path):
        """Test PPTX with complete metadata."""
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.core_properties.title = "Test Presentation"
        prs.core_properties.subject = "Unit Testing"
        prs.core_properties.author = "Test Author"
        prs.save(str(pptx_path))

        extractor = ContextExtractor(pptx_path)
        context = extractor._extract_document_context()

        assert "Title: Test Presentation" in context
        assert "Subject: Unit Testing" in context
        assert "Author: Test Author" in context


class TestSectionContextExtraction:
    """Test section-level context extraction."""

    def test_docx_with_heading(self, tmp_path):
        """Test DOCX section context with heading."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()

        # Add heading
        doc.add_heading("Introduction", level=1)

        # Add some paragraphs
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")

        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)

        # Create metadata for image at paragraph 2
        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 2}
        )

        section_context = extractor._extract_section_context(
            metadata
        )

        assert section_context == "Introduction"

    def test_docx_without_heading(self, tmp_path):
        """Test DOCX section context without heading."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")
        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 1}
        )

        section_context = extractor._extract_section_context(
            metadata
        )

        assert section_context is None

    def test_pptx_section_context(self, tmp_path):
        """Test PPTX section context (uses slide title)."""
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        extractor = ContextExtractor(pptx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={
                "slide_index": 0,
                "slide_title": "Introduction Slide"
            }
        )

        section_context = extractor._extract_section_context(
            metadata
        )

        assert section_context == "Introduction Slide"


class TestPageContextExtraction:
    """Test page-level context extraction."""

    def test_docx_page_context(self, sample_docx_path):
        """Test DOCX page context (should be None)."""
        extractor = ContextExtractor(sample_docx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 0}
        )

        page_context = extractor._extract_page_context(metadata)

        assert page_context is None

    def test_pptx_page_context(self, sample_pptx_path):
        """Test PPTX page context with slide title."""
        extractor = ContextExtractor(sample_pptx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={
                "slide_index": 0,
                "slide_title": "Test Slide"
            }
        )

        page_context = extractor._extract_page_context(metadata)

        assert page_context == "Slide: Test Slide"


class TestLocalContextExtraction:
    """Test local context extraction."""

    def test_docx_local_context(self, tmp_path):
        """Test DOCX local context extraction."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()

        # Add paragraphs
        doc.add_paragraph("Paragraph 1")
        doc.add_paragraph("Paragraph 2")
        doc.add_paragraph("Paragraph 3 (image here)")
        doc.add_paragraph("Paragraph 4")
        doc.add_paragraph("Paragraph 5")

        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 2}
        )

        local_context = extractor._extract_local_context(metadata)

        # Should include 2 before and 2 after
        assert "Paragraph 1" in local_context
        assert "Paragraph 2" in local_context
        assert "Paragraph 4" in local_context
        assert "Paragraph 5" in local_context
        # Should NOT include image paragraph
        assert "Paragraph 3 (image here)" not in local_context

    def test_docx_local_context_at_start(self, tmp_path):
        """Test DOCX local context at document start."""
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()

        doc.add_paragraph("Paragraph 1 (image here)")
        doc.add_paragraph("Paragraph 2")
        doc.add_paragraph("Paragraph 3")

        doc.save(str(docx_path))

        extractor = ContextExtractor(docx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 0}
        )

        local_context = extractor._extract_local_context(metadata)

        # Should only have paragraphs after
        assert "Paragraph 2" in local_context
        assert "Paragraph 3" in local_context

    def test_pptx_local_context(self, tmp_path):
        """Test PPTX local context extraction."""
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

        # Add text boxes
        txBox = slide.shapes.add_textbox(
            100, 100, 500, 300
        )
        txBox.text = "This is text on the slide."

        prs.save(str(pptx_path))

        extractor = ContextExtractor(pptx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"slide_index": 0}
        )

        local_context = extractor._extract_local_context(metadata)

        assert "This is text on the slide." in local_context


class TestCompleteContextExtraction:
    """Test complete context extraction workflow."""

    def test_extract_context_for_image_docx(self, tmp_path):
        """Test complete context extraction for DOCX image."""
        # Create test document
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.core_properties.title = "Test Doc"
        doc.add_heading("Section 1", level=1)
        doc.add_paragraph("Before image")
        doc.add_paragraph("Image paragraph")
        doc.add_paragraph("After image")
        doc.save(str(docx_path))

        # Create external context
        context_path = tmp_path / "context.txt"
        context_path.write_text("External context info")

        extractor = ContextExtractor(
            docx_path,
            external_context_path=context_path
        )

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={"paragraph_index": 2}
        )

        context_data = extractor.extract_context_for_image(
            metadata
        )

        assert isinstance(context_data, ContextData)
        assert context_data.external_context == \
            "External context info"
        assert "Test Doc" in context_data.document_context
        assert context_data.section_context == "Section 1"
        assert context_data.page_context is None
        assert "Before image" in context_data.local_context

    def test_extract_context_for_image_pptx(self, tmp_path):
        """Test complete context extraction for PPTX image."""
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.core_properties.title = "Test Presentation"
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
        slide.shapes.title.text = "Slide Title"
        prs.save(str(pptx_path))

        extractor = ContextExtractor(pptx_path)

        metadata = ImageMetadata(
            image_id="test_img",
            filename="test.png",
            format="PNG",
            size_bytes=1000,
            width_pixels=100,
            height_pixels=100,
            page_number=1,
            position={
                "slide_index": 0,
                "slide_title": "Slide Title"
            }
        )

        context_data = extractor.extract_context_for_image(
            metadata
        )

        assert isinstance(context_data, ContextData)
        assert "Test Presentation" in \
            context_data.document_context
        assert context_data.section_context == "Slide Title"
        assert context_data.page_context == "Slide: Slide Title"


class TestContextMerging:
    """Test context merging functionality."""

    def test_merged_context_all_levels(self):
        """Test merging with all context levels present."""
        context_data = ContextData(
            external_context="External info",
            document_context="Document: Test",
            section_context="Section 1",
            page_context="Slide: Title",
            local_context="Local text here"
        )

        merged = context_data.get_merged_context()

        assert "[External Context] External info" in merged
        assert "[Document: Document: Test]" in merged
        assert "[Section: Section 1]" in merged
        assert "[Page: Slide: Title]" in merged
        assert "[Local: Local text here]" in merged

    def test_merged_context_truncation(self):
        """Test context truncation when too long."""
        long_text = "x" * 15000

        context_data = ContextData(
            external_context=long_text,
            document_context="Test",
            local_context="Local"
        )

        merged = context_data.get_merged_context(max_chars=1000)

        assert len(merged) <= 1003  # 1000 + "..."
        assert merged.endswith("...")


# Fixtures
@pytest.fixture
def sample_docx_path(tmp_path):
    """Create a sample DOCX file for testing."""
    docx_path = tmp_path / "sample.docx"
    doc = DocxDocument()
    doc.add_paragraph("Sample paragraph")
    doc.save(str(docx_path))
    return docx_path


@pytest.fixture
def sample_pptx_path(tmp_path):
    """Create a sample PPTX file for testing."""
    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def sample_context_path():
    """Get path to sample context file."""
    return (
        Path(__file__).parent.parent /
        "fixtures" /
        "context" /
        "sample_context.txt"
    )
